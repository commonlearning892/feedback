import pandas as pd
import numpy as np
from collections import Counter
import json
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE

# Read multiple CSV files: auto-discover in the script directory
base_dir = os.path.dirname(__file__)
csv_files = []
try:
    for fname in os.listdir(base_dir):
        if fname.lower().endswith('.csv') and 'parent feedback form' in fname.lower():
            csv_files.append(os.path.join(base_dir, fname))
except Exception:
    pass

# Discover Excel files as well
xlsx_files = []
try:
    for fname in os.listdir(base_dir):
        if fname.lower().endswith('.xlsx') and 'parent feedback form' in fname.lower():
            xlsx_files.append(os.path.join(base_dir, fname))
except Exception:
    pass

# Prefer XLSX if present; otherwise fall back to CSV
input_files_used = xlsx_files if xlsx_files else csv_files

def _normalize_text(s):
    try:
        t = str(s).strip().lower()
    except Exception:
        t = str(s).lower()
    t = re.sub(r"\s+", " ", t)
    return t

def try_read_excel_with_header_detection(path):
    try:
        df0 = pd.read_excel(path, header=None)
    except Exception as e:
        raise e
    header_row = None
    candidates = ['name of the branch', 'branch', 'கிளை', 'கிளையின்']
    max_nonempty = -1
    for r in range(min(10, len(df0))):
        row_vals = [str(v).lower() if not pd.isna(v) else '' for v in list(df0.iloc[r].values)]
        hit = any(any(c in v for c in candidates) for v in row_vals)
        if hit:
            nonempty = sum(1 for v in row_vals if str(v).strip() != '')
            if nonempty > max_nonempty:
                max_nonempty = nonempty
                header_row = r
    if header_row is None:
        df = pd.read_excel(path)
    else:
        df = pd.read_excel(path, header=header_row)
    df.columns = [str(c).strip() for c in df.columns]
    return df

frames = []
for path in input_files_used:
    try:
        if os.path.exists(path):
            seg = 'Pre Primary' if 'Pre Primary' in path else ('Primary' if 'Primary' in path and 'High School' not in path else ('High School' if 'High School' in path else 'Unknown'))
            if path.lower().endswith('.xlsx'):
                try:
                    _df = try_read_excel_with_header_detection(path)
                except Exception as e:
                    print(f"Warning: could not read Excel {path}: {e}. If missing, install the 'openpyxl' package.")
                    continue
            elif path.lower().endswith('.csv'):
                _df = pd.read_csv(path)
            else:
                continue
            try:
                _df.columns = [str(c).strip() for c in _df.columns]
            except Exception:
                pass
            _df['Segment'] = seg
            frames.append(_df)
    except Exception as e:
        print(f"Warning: could not read {path}: {e}")

# If XLSX files were preferred but none could be read, fall back to CSVs if available
if not frames and xlsx_files and csv_files:
    print("Note: Could not parse Excel files; falling back to CSV. To use Excel inputs, install 'openpyxl'.")
    for path in csv_files:
        try:
            if os.path.exists(path):
                seg = 'Pre Primary' if 'Pre Primary' in path else ('Primary' if 'Primary' in path and 'High School' not in path else ('High School' if 'High School' in path else 'Unknown'))
                _df = pd.read_csv(path)
                _df['Segment'] = seg
                frames.append(_df)
        except Exception as e:
            print(f"Warning: could not read {path}: {e}")

if not frames:
    raise FileNotFoundError("No input CSV/XLSX files found. Expected Pre Primary, Primary, or High School files in the folder. If using Excel, ensure 'openpyxl' is installed.")

df = pd.concat(frames, ignore_index=True)

# Remove empty rows
df = df.dropna(how='all')

# Print basic info
print(f"Total records: {len(df)}")
print(f"\nColumns ({len(df.columns)}):")
for i, col in enumerate(df.columns, 1):
    print(f"{i}. {col}")

# Define rating mapping for conversion
rating_map = {
    'Excellent(மிகநன்று)': 5,
    'Excellent(மிக நன்று)': 5,
    'Good(நன்று)': 4,
    'Average(சராசரி)': 3,
    'Satisfactory(சராசரி)': 3,
    'Satisfactory(திருப்தி)': 3,
    'Need Improvement(முன்னேற்றம் தேவை)': 2,
    'Needs Improvement(முன்னேற்றம் தேவை)': 2,
    'Needs to Improve(முன்னேற்றம் தேவை)': 2,
    'Poor(மோசம்)': 1,
    'Not Applicable(பொருந்தாது)': 0,
    'Not Applicable(பொருந்தாது': 0,
    'Yes(ஆம்)': 5,
    'Maybe(இருக்கலாம்)': 3,
    'No(இல்லை)': 1
}

def canonicalize_rating(value):
    if pd.isna(value) or value == '':
        return None
    s = str(value).strip().lower()
    # Tamil/localized variants included
    if 'not applicable' in s or 'பொருந்தாது' in s:
        return 'Not Applicable'
    if 'excellent' in s or 'மிகநன்று' in s or 'மிக நன்று' in s:
        return 'Excellent'
    if 'good' in s or 'நன்று' in s:
        return 'Good'
    if 'average' in s or 'satisfactory' in s or 'திருப்தி' in s or 'சராசரி' in s:
        return 'Average'
    if 'poor' in s or 'மோசம்' in s or 'needs' in s or 'need' in s or 'improve' in s or 'முன்னேற்றம்' in s:
        return 'Poor'
    return None

def normalize_rating(value):
    """Convert rating strings to numeric values"""
    if pd.isna(value) or value == '':
        return 0
    value_str = str(value).strip()
    if value_str in rating_map:
        return rating_map[value_str]
    # Fallback: canonicalize and map
    cat = canonicalize_rating(value_str)
    if cat == 'Excellent':
        return 5
    if cat == 'Good':
        return 4
    if cat == 'Average':
        return 3
    if cat == 'Poor':
        return 1
    if cat == 'Not Applicable':
        return 0
    return 0

# Extract key columns
branch_col = 'Name of the Branch( கிளையின் பெயர்)'
class_col = 'Class( வகுப்பு )'
orientation_col = 'Orientation( பயிற்சி வகை )'
language_col = 'II Language( இரண்டாம் மொழிப்பாடம்) '

def normalize_branch_name(s):
    if pd.isna(s):
        return 'Unknown'
    t = str(s).strip()
    t = re.sub(r'\s+', ' ', t)
    t = re.sub(r'[\-_/]+$', '', t).strip()
    t = t.title()
    return t

def branch_canonical_key(s):
    if pd.isna(s):
        return 'unknown'
    t = str(s).lower().strip()
    t = re.sub(r'\b(sri|chaitanya|techno|technos|school|schools|campus|branch)\b', '', t)
    t = re.sub(r'[^a-z0-9]', '', t)
    return t

# Resolve key columns by keyword matching; create if missing
def resolve_column_name(df_in, expected, keywords):
    try:
        cols = list(df_in.columns)
    except Exception:
        cols = []
    if expected in cols:
        return expected
    best = None
    best_score = 0
    for c in cols:
        low = _normalize_text(c)
        score = sum(1 for kw in keywords if kw in low)
        if score > best_score:
            best_score = score
            best = c
    return best if best else expected

def resolve_or_create(df_in, expected, keywords):
    col = resolve_column_name(df_in, expected, keywords)
    if col not in df_in.columns:
        df_in[col] = None
    return col

branch_col = resolve_or_create(df, branch_col, ['name of the branch', 'branch', 'கிளை'])
class_col = resolve_or_create(df, class_col, ['class', 'வகுப்பு'])
orientation_col = resolve_or_create(df, orientation_col, ['orientation', 'பயிற்சி'])
language_col = resolve_or_create(df, language_col, ['ii language', 'language', 'இரண்டாம்', 'மொழி'])

# Normalize branch names
df[branch_col] = df[branch_col].apply(normalize_branch_name)
# Club similar branch names: compute canonical key and remap to most frequent display variant
try:
    df['BranchKey'] = df[branch_col].apply(branch_canonical_key)
    display_map = {}
    for key, grp in df.groupby('BranchKey'):
        name_counts = grp[branch_col].value_counts()
        # choose most frequent; tie-breaker: shortest name
        top_count = name_counts.max()
        candidates = [n for n, c in name_counts.items() if c == top_count]
        display = sorted(candidates, key=lambda x: (len(str(x)), str(x)))[0]
        display_map[key] = display
    df[branch_col] = df['BranchKey'].map(display_map).fillna(df[branch_col])
    df.drop(columns=['BranchKey'], inplace=True)
except Exception:
    pass

df[class_col] = df[class_col].fillna('Unknown').str.strip()
df[orientation_col] = df[orientation_col].fillna('Unknown').str.strip()
df[language_col] = df[language_col].fillna('Unknown').str.strip()

# Subject feedback columns (dynamic detection across segments)
subject_cols = {}
for col in df.columns:
    lowc = str(col).lower()
    if 'subject wise feedback' in lowc:
        # Extract name inside square brackets if present
        m = re.search(r'\[(.*?)\]', str(col))
        if m:
            name = m.group(1).strip()
        else:
            # Fallback: take text after the keyword
            name = str(col).split(')', 1)[-1].strip()[:60]
        # Shorten noisy names
        name = re.sub(r'\s+', ' ', name)
        name = name.replace('skills', '').replace(' - ', ' ').strip()
        # Avoid duplicates
        base = name if name else 'Subject'
        key = base
        i = 2
        while key in subject_cols:
            key = f"{base} {i}"
            i += 1
        subject_cols[key] = col

# Environment quality columns (shortened names for analysis)
env_cols = {}
infra_cols = {}
parent_cols = {}
admin_cols = {}

for col in df.columns:
    lowc = str(col).lower()
    if 'overall quality of the school environment' in lowc:
        match = re.search(r'\[(.*?)\]', col)
        if match:
            key = match.group(1).split('(')[0].strip()[:50]
            env_cols[key] = col
    elif 'overall school infrastructure' in lowc:
        match = re.search(r'\[(.*?)\]', col)
        if match:
            key = match.group(1).split('(')[0].strip()[:50]
            infra_cols[key] = col
    elif ('parent-teacher' in lowc or 'parent–teacher' in lowc or 'ptm' in lowc):
        match = re.search(r'\[(.*?)\]', col)
        if match:
            key = match.group(1).split('(')[0].strip()[:50]
            parent_cols[key] = col
    elif ('administration' in lowc or 'admin team' in lowc or 'front office' in lowc or 'leadership' in lowc or 'principal' in lowc or 'vice principal' in lowc or 'coordinator' in lowc or 'administrative support' in lowc):
        match = re.search(r'\[(.*?)\]', col)
        if match:
            key = match.group(1).split('(')[0].strip()[:50]
            admin_cols[key] = col

# Convert ratings to numeric
for col in df.columns:
    if col not in [branch_col, class_col, orientation_col, language_col, 'Timestamp', 
                   'Student Name( மாணவர் பெயர்)', 'SCS NUMBER( SCS எண்)',
                   'Parent Name( பெற்றோர் பெயர்) ', 'Parent Phone Number( பெற்றோர் தொலைபேசி எண்)']:
        df[col + '_numeric'] = df[col].apply(normalize_rating)

# Calculate overall scores
def existing_numeric(cols_dict):
    return [col + '_numeric' for col in cols_dict.values() if (col + '_numeric') in df.columns]

def rowwise_mean_from_cols(cols_dict):
    series_list = []
    for col in cols_dict.values():
        if col in df.columns:
            num_col = col + '_numeric'
            if num_col in df.columns:
                s = df[num_col]
            else:
                s = df[col].apply(normalize_rating)
            series_list.append(s.replace(0, np.nan))
    if not series_list:
        return np.nan
    return pd.concat(series_list, axis=1).mean(axis=1)

subject_numeric_cols = existing_numeric(subject_cols)
env_numeric_cols = existing_numeric(env_cols)
infra_numeric_cols = existing_numeric(infra_cols)
parent_numeric_cols = existing_numeric(parent_cols)
admin_numeric_cols = existing_numeric(admin_cols)

# Calculate average scores per student
df['Subject_Avg'] = rowwise_mean_from_cols(subject_cols)
df['Environment_Avg'] = rowwise_mean_from_cols(env_cols)
df['Infrastructure_Avg'] = rowwise_mean_from_cols(infra_cols)
df['Parent_Teacher_Avg'] = rowwise_mean_from_cols(parent_cols)
df['Admin_Avg'] = rowwise_mean_from_cols(admin_cols)
df['Overall_Avg'] = df[[col for col in df.columns if col.endswith('_numeric')]].replace(0, np.nan).mean(axis=1)

# Additional column detections for advanced dashboards
def find_columns(keyword_substrings):
    cols = []
    for col in df.columns:
        low = str(col).lower()
        if all(kw in low for kw in keyword_substrings):
            cols.append(col)
    return cols

recommend_cols = [c for c in df.columns if 'recommend' in str(c).lower()]
app_cols = [c for c in df.columns if ' app' in str(c).lower() or 'app ' in str(c).lower() or 'application' in str(c).lower()]
transport_cols = [c for c in df.columns if 'transport' in str(c).lower()]
ptm_cols = [c for c in df.columns if 'parent-teacher' in str(c).lower() or 'parent–teacher' in str(c).lower() or 'ptm' in str(c).lower()]

sat_cols = [c for c in df.columns if 'mention two areas' in str(c).lower() or 'most satisfied' in str(c).lower() or ('you are most satisfied' in str(c).lower())]
improve_cols = [c for c in df.columns if 'further improve' in str(c).lower() or 'can further improve' in str(c).lower() or ('choose two areas' in str(c).lower() and 'improve' in str(c).lower())]

# Teaching indicators
clarity_cols = [c for c in df.columns if 'concept' in str(c).lower() or 'clarity' in str(c).lower()]
approach_cols = [c for c in df.columns if 'approach' in str(c).lower() or 'approachability' in str(c).lower()]
engagement_cols = [c for c in df.columns if 'engage' in str(c).lower() or 'stories' in str(c).lower() or 'rhymes' in str(c).lower() or 'activities' in str(c).lower()]
comm_skill_cols = [c for c in df.columns if 'communication skills' in str(c).lower()]

# Concern handling (role-wise)
concern_roles_map = {}
for col in df.columns:
    low = str(col).lower()
    if 'addresses my concerns' in low or 'handles my concerns' in low or 'addresses concerns' in low:
        m = re.search(r'\[(.*?)\]', str(col))
        if m:
            role = m.group(1).strip()
            concern_roles_map[role] = col

# Concern resolution satisfaction (Yes/No/NA)
concern_resolve_cols = [c for c in df.columns if ('handles concerns' in str(c).lower() or 'actions taken' in str(c).lower() or 'resolve' in str(c).lower()) and ('select' in str(c).lower() or 'not applicable' in str(c).lower())]

def mean_cols(cols):
    if not cols:
        return None
    num_cols = [c + '_numeric' for c in cols if c + '_numeric' in df.columns]
    if not num_cols:
        return None
    return float(df[num_cols].replace(0, np.nan).mean(axis=1).mean())

# Canonicalize rating label to one of: Excellent, Good, Average, Poor, Not Applicable
def canonicalize_rating(value):
    if pd.isna(value) or value == '':
        return None
    s = str(value).strip().lower()
    # Tamil/localized variants included
    if 'not applicable' in s or 'பொருந்தாது' in s:
        return 'Not Applicable'
    if 'excellent' in s or 'மிகநன்று' in s or 'மிக நன்று' in s:
        return 'Excellent'
    if 'good' in s or 'நன்று' in s:
        return 'Good'
    if 'average' in s or 'satisfactory' in s or 'திருப்தி' in s or 'சராசரி' in s:
        return 'Average'
    if 'poor' in s or 'மோசம்' in s or 'needs' in s or 'need' in s or 'improve' in s or 'முன்னேற்றம்' in s:
        return 'Poor'
    return None

def parse_reasons(value):
    if pd.isna(value) or value == '':
        return []
    parts = re.split(r'[;|\n]+', str(value))
    return [p.strip() for p in parts if p and p.strip()]

def bucket_reason(text):
    s = str(text).lower()
    if 'no concerns' in s or 'மேற்கண்ட எதுவுமில்லை' in s:
        return 'No Concerns'
    if 'transport' in s:
        return 'Transport'
    if 'sports' in s:
        return 'Sports'
    if 'infrastructure' in s or 'facilities' in s or 'facility' in s:
        return 'Infrastructure & Facilities'
    if 'discipline' in s or 'values' in s:
        return 'Discipline & Values'
    if 'competition' in s or 'event' in s or 'celebration' in s:
        return 'Events & Celebrations'
    if 'student communication' in s or ('communication' in s and 'student' in s):
        return 'Student Communication'
    if 'communication' in s:
        return 'Communication'
    if 'app' in s:
        return 'App'
    if 'academics' in s or 'teaching' in s or 'stories' in s or 'rhymes' in s or 'activities' in s:
        return 'Academics & Activities'
    if 'environment' in s or 'hygiene' in s or 'clean' in s or 'safety' in s:
        return 'Environment & Safety'
    return 'Other'

# Classify Yes/No/Maybe/Not Applicable for localized strings
def classify_ynm(value):
    if pd.isna(value) or value == '':
        return None
    s = str(value).strip().lower()
    if ('yes' in s) or ('ஆம்' in s):
        return 'Yes'
    if ('no' in s) or ('இல்லை' in s):
        return 'No'
    if ('maybe' in s) or ('இருக்கலாம்' in s):
        return 'Maybe'
    if ('not applicable' in s) or ('பொருந்தாது' in s):
        return 'Not Applicable'
    return None

# Create comprehensive statistics
stats = {
    'summary': {
        'total_responses': len(df),
        'branches': df[branch_col].value_counts().to_dict(),
        'classes': df[class_col].value_counts().to_dict(),
        'orientations': df[orientation_col].value_counts().to_dict(),
        'languages': df[language_col].value_counts().to_dict()
    },
    'branch_performance': {},
    'orientation_performance': {},
    'class_performance': {},
    'subject_performance': {},
    'category_performance': {
        'Environment Quality': {},
        'Infrastructure': {},
        'Parent-Teacher Interaction': {},
        'Administrative Support': {}
    }
}

# Branch-wise analysis
for branch in df[branch_col].unique():
    if branch and branch != 'Unknown' and branch != '':
        branch_data = df[df[branch_col] == branch]
        stats['branch_performance'][branch] = {
            'count': len(branch_data),
            'subject_avg': float(branch_data['Subject_Avg'].mean()),
            'environment_avg': float(branch_data['Environment_Avg'].mean()),
            'infrastructure_avg': float(branch_data['Infrastructure_Avg'].mean()),
            'parent_teacher_avg': float(branch_data['Parent_Teacher_Avg'].mean()),
            'admin_avg': float(branch_data['Admin_Avg'].mean()),
            'overall_avg': float(branch_data['Overall_Avg'].mean())
        }

# Orientation-wise analysis
for orientation in df[orientation_col].unique():
    if orientation and orientation != 'Unknown' and orientation != '':
        orient_data = df[df[orientation_col] == orientation]
        stats['orientation_performance'][orientation] = {
            'count': len(orient_data),
            'subject_avg': float(orient_data['Subject_Avg'].mean()),
            'environment_avg': float(orient_data['Environment_Avg'].mean()),
            'infrastructure_avg': float(orient_data['Infrastructure_Avg'].mean()),
            'parent_teacher_avg': float(orient_data['Parent_Teacher_Avg'].mean()),
            'admin_avg': float(orient_data['Admin_Avg'].mean()),
            'overall_avg': float(orient_data['Overall_Avg'].mean())
        }

# Class-wise analysis
for class_name in df[class_col].unique():
    if class_name and class_name != 'Unknown' and class_name != '':
        class_data = df[df[class_col] == class_name]
        stats['class_performance'][class_name] = {
            'count': len(class_data),
            'subject_avg': float(class_data['Subject_Avg'].mean()),
            'environment_avg': float(class_data['Environment_Avg'].mean()),
            'infrastructure_avg': float(class_data['Infrastructure_Avg'].mean()),
            'parent_teacher_avg': float(class_data['Parent_Teacher_Avg'].mean()),
            'admin_avg': float(class_data['Admin_Avg'].mean()),
            'overall_avg': float(class_data['Overall_Avg'].mean())
        }

# Subject-wise analysis (overall)
for subject_name, subject_col in subject_cols.items():
    if subject_col not in df.columns:
        continue
    num_col = subject_col + '_numeric'
    if num_col in df.columns:
        ratings = df[num_col].replace(0, np.nan)
    else:
        ratings = df[subject_col].apply(normalize_rating).replace(0, np.nan)
    stats['subject_performance'][subject_name] = {
        'average': float(ratings.mean()),
        'excellent_count': int((df[subject_col].astype(str).str.lower().str.contains('excellent')).sum()),
        'good_count': int((df[subject_col].astype(str).str.lower().str.contains('good')).sum()),
        'average_count': int((df[subject_col].astype(str).str.lower().str.contains('average|satisfactory|திருப்தி|சராசரி')).sum()),
        'poor_count': int((df[subject_col].astype(str).str.lower().str.contains('poor|need|needs|improve|முன்னேற்றம்|மோசம்')).sum()),
        'rating_distribution': df[subject_col].value_counts().to_dict()
    }

# Per-branch subject-wise analysis
branch_subject_perf = {}
for branch, group in df.groupby(branch_col):
    subd = {}
    for subject_name, subject_col in subject_cols.items():
        if subject_col not in group.columns:
            continue
        num_col = subject_col + '_numeric'
        if num_col in group.columns:
            ratings = group[num_col].replace(0, np.nan)
        else:
            ratings = group[subject_col].apply(normalize_rating).replace(0, np.nan)
        dist = group[subject_col].value_counts().to_dict()
        subd[subject_name] = {
            'average': float(ratings.mean()),
            'rating_distribution': dist
        }
    branch_subject_perf[branch] = subd
stats['branch_subject_performance'] = branch_subject_perf

# Global per-segment subject-wise analysis
segment_subject_perf = {}
for seg, gseg in df.groupby('Segment'):
    subd = {}
    for subject_name, subject_col in subject_cols.items():
        if subject_col not in gseg.columns:
            continue
        dist = gseg[subject_col].dropna().value_counts().to_dict()
        total = sum(dist.values())
        if total == 0:
            continue
        num_col = subject_col + '_numeric'
        if num_col in gseg.columns:
            ratings = gseg[num_col].replace(0, np.nan)
        else:
            ratings = gseg[subject_col].apply(normalize_rating).replace(0, np.nan)
        subd[subject_name] = {
            'average': float(ratings.mean()),
            'rating_distribution': dist
        }
    if subd:
        segment_subject_perf[seg] = subd
stats['segment_subject_performance'] = segment_subject_perf

# Per-branch, per-segment, subject-wise analysis
branch_segment_subject_perf = {}
for branch, g_branch in df.groupby(branch_col):
    branch_segment_subject_perf[branch] = {}
    for seg, gseg in g_branch.groupby('Segment'):
        subd = {}
        for subject_name, subject_col in subject_cols.items():
            if subject_col not in gseg.columns:
                continue
            dist = gseg[subject_col].dropna().value_counts().to_dict()
            total = sum(dist.values())
            if total == 0:
                continue
            num_col = subject_col + '_numeric'
            if num_col in gseg.columns:
                ratings = gseg[num_col].replace(0, np.nan)
            else:
                ratings = gseg[subject_col].apply(normalize_rating).replace(0, np.nan)
            subd[subject_name] = {
                'average': float(ratings.mean()),
                'rating_distribution': dist
            }
        if subd:
            branch_segment_subject_perf[branch][seg] = subd
stats['branch_segment_subject_performance'] = branch_segment_subject_perf

# Environment quality detailed analysis (overall)
for env_name, env_col in env_cols.items():
    env_numeric = env_col + '_numeric'
    ratings = df[env_numeric].replace(0, np.nan)
    stats['category_performance']['Environment Quality'][env_name] = {
        'average': float(ratings.mean()),
        'rating_distribution': df[env_col].value_counts().to_dict()
    }

# Infrastructure detailed analysis (overall)
for infra_name, infra_col in infra_cols.items():
    infra_numeric = infra_col + '_numeric'
    ratings = df[infra_numeric].replace(0, np.nan)
    stats['category_performance']['Infrastructure'][infra_name] = {
        'average': float(ratings.mean()),
        'rating_distribution': df[infra_col].value_counts().to_dict()
    }

# Parent-Teacher interaction analysis (overall)
for parent_name, parent_col in parent_cols.items():
    parent_numeric = parent_col + '_numeric'
    ratings = df[parent_numeric].replace(0, np.nan)
    stats['category_performance']['Parent-Teacher Interaction'][parent_name] = {
        'average': float(ratings.mean()),
        'rating_distribution': df[parent_col].value_counts().to_dict()
    }

# Administrative support analysis (overall)
for admin_name, admin_col in admin_cols.items():
    admin_numeric = admin_col + '_numeric'
    ratings = df[admin_numeric].replace(0, np.nan)
    stats['category_performance']['Administrative Support'][admin_name] = {
        'average': float(ratings.mean()),
        'rating_distribution': df[admin_col].value_counts().to_dict()
    }

# Rankings - Top performers
def create_ranking(data_dict, metric='overall_avg'):
    ranking = sorted(data_dict.items(), key=lambda x: x[1].get(metric, 0), reverse=True)
    return [(name, data[metric], data['count']) for name, data in ranking if not np.isnan(data.get(metric, 0))]

stats['rankings'] = {
    'branches': create_ranking(stats['branch_performance']),
    'orientations': create_ranking(stats['orientation_performance']),
    'classes': create_ranking(stats['class_performance']),
    'subjects': [(name, data['average']) for name, data in 
                 sorted(stats['subject_performance'].items(), key=lambda x: x[1]['average'], reverse=True)]
}

# Replace NaN values with None (null in JSON)
def clean_nan(obj):
    if isinstance(obj, dict):
        return {k: clean_nan(v) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [clean_nan(item) for item in obj]
    elif isinstance(obj, float) and np.isnan(obj):
        return None
    return obj


# Executive summary KPIs and additional aggregations
stats['summary']['overall_avg'] = float(df['Overall_Avg'].mean())

category_scores = {
    'Academics': float(df['Subject_Avg'].mean()),
    'Administration': float(df['Admin_Avg'].mean()),
    'Environment': float(df['Environment_Avg'].mean()),
    'Infrastructure': float(df['Infrastructure_Avg'].mean()),
}
app_avg = mean_cols(app_cols)
transport_avg = mean_cols(transport_cols)
if app_avg is not None:
    category_scores['App'] = app_avg
if transport_avg is not None:
    category_scores['Transport'] = transport_avg
stats['summary']['category_scores'] = category_scores

# Recommendation distribution and percent Yes (robust to localized values)
rec_counts = {'Yes': 0, 'No': 0, 'Maybe': 0, 'Not Applicable': 0}
if recommend_cols:
    for col in recommend_cols:
        mapped = df[col].apply(classify_ynm)
        vc = mapped.value_counts()
        for k, v in vc.items():
            if k in rec_counts:
                rec_counts[k] += int(v)
total_rec = rec_counts['Yes'] + rec_counts['No'] + rec_counts['Maybe']
yes_pct = (rec_counts['Yes'] / total_rec * 100.0) if total_rec > 0 else None
stats['recommendation'] = {
    'distribution': rec_counts,
    'yes_pct': yes_pct
}

# PTM effectiveness (overall and per branch)
ptm_avg = mean_cols(ptm_cols)
stats['ptm_effectiveness'] = ptm_avg
ptm_by_branch = {}
for branch, group in df.groupby(branch_col):
    if not ptm_cols:
        ptm_by_branch[branch] = None
    else:
        num_cols = [c + '_numeric' for c in ptm_cols if c + '_numeric' in group.columns]
        if num_cols:
            ptm_by_branch[branch] = float(group[num_cols].replace(0, np.nan).mean(axis=1).mean())
        else:
            ptm_by_branch[branch] = None
stats['ptm_effectiveness_by_branch'] = ptm_by_branch

# Teaching indicators aggregation
stats['teaching_indicators'] = {
    'Concept Clarity': mean_cols(clarity_cols),
    'Teacher Approachability': mean_cols(approach_cols),
    'Engagement': mean_cols(engagement_cols),
    'Communication Skills': mean_cols(comm_skill_cols),
}

# Environment focus metrics (pick common keywords from Environment Quality group)
def compute_env_focus(df_subset):
    env_focus_local = {}
    for key, col in env_cols.items():
        if (col + '_numeric') not in df_subset.columns:
            continue
        low = key.lower()
        avg = float(df_subset[col + '_numeric'].replace(0, np.nan).mean())
        if 'interest' in low or 'enthusiasm' in low:
            env_focus_local['Interest in attending school'] = avg
        elif 'secure' in low or 'safety' in low:
            env_focus_local['Campus safety'] = avg
        elif 'moral' in low or 'values' in low:
            env_focus_local['Moral values'] = avg
        elif 'social' in low or 'confidence' in low:
            env_focus_local['Social confidence'] = avg
    return env_focus_local

env_focus = compute_env_focus(df)
stats['environment_focus'] = env_focus

env_focus_by_branch = {}
for branch, group in df.groupby(branch_col):
    env_focus_by_branch[branch] = compute_env_focus(group)
stats['environment_focus_by_branch'] = env_focus_by_branch

# Communication metrics (common keywords) overall and per branch
def compute_comm_metrics(df_subset):
    cm = {}
    for col in df.columns:
        low = str(col).lower()
        label = None
        if 'front office' in low or 'front-office' in low:
            label = 'Front Office Support'
        elif 'leadership' in low or 'principal access' in low or 'access' in low:
            label = 'Leadership Access'
        elif 'app' in low and ('usability' in low or 'use' in low):
            label = 'App Usability'
        elif 'timely updates' in low or 'timely' in low or 'updates' in low:
            label = 'Timely Updates'
        if label:
            num_col = col + '_numeric'
            if num_col in df_subset.columns:
                cm[label] = float(df_subset[num_col].replace(0, np.nan).mean())
    return cm

comm_metrics = compute_comm_metrics(df)
stats['communication_metrics'] = comm_metrics

comm_by_branch = {}
for branch, group in df.groupby(branch_col):
    comm_by_branch[branch] = compute_comm_metrics(group)
stats['communication_metrics_by_branch'] = comm_by_branch

# Concern handling role-wise averages (overall and per branch)
concern_roles = {}
for role, col in concern_roles_map.items():
    num_col = col + '_numeric'
    if num_col in df.columns:
        concern_roles[role] = float(df[num_col].replace(0, np.nan).mean())
stats['concern_roles'] = concern_roles

concern_roles_by_branch = {}
for branch, group in df.groupby(branch_col):
    vals = {}
    for role, col in concern_roles_map.items():
        num_col = col + '_numeric'
        if num_col in group.columns:
            vals[role] = float(group[num_col].replace(0, np.nan).mean())
    concern_roles_by_branch[branch] = vals
stats['concern_roles_by_branch'] = concern_roles_by_branch

# Concern resolution distribution (Yes/No/Not Applicable)
concern_dist = {'Yes': 0, 'No': 0, 'Not Applicable': 0, 'Maybe': 0}
for col in concern_resolve_cols:
    mapped = df[col].apply(classify_ynm)
    vc = mapped.value_counts()
    for k, v in vc.items():
        if k in concern_dist:
            concern_dist[k] += int(v)
stats['concern_resolution'] = concern_dist

# Branch recommendation percentage
branch_rec_pct = {}
if recommend_cols:
    for branch, group in df.groupby(branch_col):
        yes = 0
        total = 0
        for col in recommend_cols:
            mapped = group[col].apply(classify_ynm)
            vc = mapped.value_counts()
            yes += int(vc.get('Yes', 0))
            total += int(vc.get('Yes', 0)) + int(vc.get('No', 0)) + int(vc.get('Maybe', 0))
        branch_rec_pct[branch] = (yes / total * 100.0) if total > 0 else None
stats['branch_recommendation_pct'] = branch_rec_pct

# Branch recommendation counts (Yes/No/Maybe/Not Applicable)
branch_rec_counts = {}
if recommend_cols:
    for branch, group in df.groupby(branch_col):
        counts = {'Yes': 0, 'No': 0, 'Maybe': 0, 'Not Applicable': 0}
        for col in recommend_cols:
            mapped = group[col].apply(classify_ynm)
            vc = mapped.value_counts()
            for k, v in vc.items():
                if k in counts:
                    counts[k] += int(v)
        branch_rec_counts[branch] = counts
stats['branch_recommendation_counts'] = branch_rec_counts

# Branch concern resolution counts
branch_concern_counts = {}
if concern_resolve_cols:
    for branch, group in df.groupby(branch_col):
        counts = {'Yes': 0, 'No': 0, 'Maybe': 0, 'Not Applicable': 0}
        for col in concern_resolve_cols:
            mapped = group[col].apply(classify_ynm)
            vc = mapped.value_counts()
            for k, v in vc.items():
                if k in counts:
                    counts[k] += int(v)
        branch_concern_counts[branch] = counts
stats['branch_concern_resolution'] = branch_concern_counts

# Branch rating category counts by group (Subjects/Environment/Infrastructure/Parent-Teacher/Admin)
def count_ratings_for_group(df_group, cols):
    counts = {'Excellent': 0, 'Good': 0, 'Average': 0, 'Poor': 0}
    for col in cols:
        if col not in df_group.columns:
            continue
        series = df_group[col].dropna()
        for val in series:
            cat = canonicalize_rating(val)
            if cat in counts:
                counts[cat] += 1
    return counts

branch_rating_counts = {}
for branch, group in df.groupby(branch_col):
    branch_rating_counts[branch] = {
        'Subjects': count_ratings_for_group(group, list(subject_cols.values())),
        'Environment': count_ratings_for_group(group, list(env_cols.values())),
        'Infrastructure': count_ratings_for_group(group, list(infra_cols.values())),
        'Parent-Teacher': count_ratings_for_group(group, list(parent_cols.values())),
        'Administrative Support': count_ratings_for_group(group, list(admin_cols.values())),
    }
stats['branch_rating_counts'] = branch_rating_counts

# Helper: compute per-branch recommendation counts for a dataframe subset
def compute_branch_rec_counts(df_subset):
    out = {}
    if not recommend_cols:
        return out
    for branch, group in df_subset.groupby(branch_col):
        counts = {'Yes': 0, 'No': 0, 'Maybe': 0, 'Not Applicable': 0}
        for col in recommend_cols:
            mapped = group[col].apply(classify_ynm)
            vc = mapped.value_counts()
            for k, v in vc.items():
                if k in counts:
                    counts[k] += int(v)
        out[branch] = counts
    return out

# Helper: compute per-branch rating counts (E/G/A/P) for each group for a dataframe subset
def compute_branch_rating_counts(df_subset):
    out = {}
    for branch, group in df_subset.groupby(branch_col):
        out[branch] = {
            'Subjects': count_ratings_for_group(group, list(subject_cols.values())),
            'Environment': count_ratings_for_group(group, list(env_cols.values())),
            'Infrastructure': count_ratings_for_group(group, list(infra_cols.values())),
            'Parent-Teacher': count_ratings_for_group(group, list(parent_cols.values())),
            'Administrative Support': count_ratings_for_group(group, list(admin_cols.values())),
        }
    return out

# Build breakdowns by class, orientation, and their pair (intersection)
classes_unique = [c for c in sorted(df[class_col].dropna().unique()) if c and c != 'Unknown']
orients_unique = [o for o in sorted(df[orientation_col].dropna().unique()) if o and o != 'Unknown']

brc_by_class = {}
for c in classes_unique:
    brc_by_class[c] = compute_branch_rec_counts(df[df[class_col] == c])
brc_by_orient = {}
for o in orients_unique:
    brc_by_orient[o] = compute_branch_rec_counts(df[df[orientation_col] == o])
brc_by_pair = {}
for c in classes_unique:
    brc_by_pair[c] = {}
    for o in orients_unique:
        sub = df[(df[class_col] == c) & (df[orientation_col] == o)]
        brc_by_pair[c][o] = compute_branch_rec_counts(sub)
stats['branch_recommendation_counts_by'] = {
    'class': brc_by_class,
    'orientation': brc_by_orient,
    'pair': brc_by_pair
}

brg_by_class = {}
for c in classes_unique:
    brg_by_class[c] = compute_branch_rating_counts(df[df[class_col] == c])
brg_by_orient = {}
for o in orients_unique:
    brg_by_orient[o] = compute_branch_rating_counts(df[df[orientation_col] == o])
brg_by_pair = {}
for c in classes_unique:
    brg_by_pair[c] = {}
    for o in orients_unique:
        sub = df[(df[class_col] == c) & (df[orientation_col] == o)]
        brg_by_pair[c][o] = compute_branch_rating_counts(sub)
stats['branch_rating_counts_by'] = {
    'class': brg_by_class,
    'orientation': brg_by_orient,
    'pair': brg_by_pair
}

rec_reasons = { 'Yes': Counter(), 'Maybe': Counter(), 'No': Counter() }
rec_col = recommend_cols[0] if recommend_cols else None
if rec_col:
    for _, row in df.iterrows():
        status = classify_ynm(row.get(rec_col))
        if status not in rec_reasons:
            continue
        # Choose strengths for Yes, improvement for Maybe/No
        cols = sat_cols if status == 'Yes' else improve_cols
        reasons = []
        for c in cols:
            reasons.extend(parse_reasons(row.get(c)))
        buckets = [bucket_reason(r) for r in reasons]
        for b in buckets:
            rec_reasons[status][b] += 1

# Raw reasons (exact selections from CSV) for transparency in UI
rec_reasons_raw = { 'Yes': Counter(), 'Maybe': Counter(), 'No': Counter() }
if rec_col:
    for _, row in df.iterrows():
        status = classify_ynm(row.get(rec_col))
        if status not in rec_reasons_raw:
            continue
        cols = sat_cols if status == 'Yes' else improve_cols
        reasons = []
        for c in cols:
            reasons.extend(parse_reasons(row.get(c)))
        for r in reasons:
            if r:
                rec_reasons_raw[status][r] += 1

def reasons_to_top(counter):
    total = sum(counter.values())
    if total == 0:
        return {'total_reasons': 0, 'top': [], 'top_detail': []}
    items = counter.most_common(8)
    return {
        'total_reasons': int(total),
        'top': [[k, round(v*100.0/total, 1)] for k, v in items],
        'top_detail': [[k, int(v), round(v*100.0/total, 1)] for k, v in items]
    }

stats['recommendation_reasons'] = {
    'Yes': reasons_to_top(rec_reasons['Yes']),
    'Maybe': reasons_to_top(rec_reasons['Maybe']),
    'No': reasons_to_top(rec_reasons['No'])
}

stats['recommendation_reasons_raw'] = {
    'Yes': reasons_to_top(rec_reasons_raw['Yes']),
    'Maybe': reasons_to_top(rec_reasons_raw['Maybe']),
    'No': reasons_to_top(rec_reasons_raw['No'])
}

# Per-branch, per-segment aggregates for side-by-side comparisons
branch_segment_perf = {}
branch_segment_rec_counts = {}
branch_segment_rec_reasons = {}
segments = [s for s in df['Segment'].dropna().unique()]
for branch, g_branch in df.groupby(branch_col):
    branch_segment_perf[branch] = {}
    branch_segment_rec_counts[branch] = {}
    branch_segment_rec_reasons[branch] = {}
    for seg, g in g_branch.groupby('Segment'):
        branch_segment_perf[branch][seg] = {
            'count': int(len(g)),
            'subject_avg': float(g['Subject_Avg'].mean()),
            'environment_avg': float(g['Environment_Avg'].mean()),
            'infrastructure_avg': float(g['Infrastructure_Avg'].mean()),
            'parent_teacher_avg': float(g['Parent_Teacher_Avg'].mean()),
            'admin_avg': float(g['Admin_Avg'].mean()),
            'overall_avg': float(g['Overall_Avg'].mean())
        }
        counts = {'Yes': 0, 'No': 0, 'Maybe': 0, 'Not Applicable': 0}
        if 'recommend_cols' in globals() and recommend_cols:
            for col in recommend_cols:
                mapped = g[col].apply(classify_ynm)
                vc = mapped.value_counts()
                for k, v in vc.items():
                    if k in counts:
                        counts[k] += int(v)
        branch_segment_rec_counts[branch][seg] = counts
        rec_col = recommend_cols[0] if ('recommend_cols' in globals() and recommend_cols) else None
        seg_reasons = { 'Yes': Counter(), 'Maybe': Counter(), 'No': Counter() }
        if rec_col:
            for _, row in g.iterrows():
                status = classify_ynm(row.get(rec_col))
                if status not in seg_reasons:
                    continue
                cols = sat_cols if status == 'Yes' else improve_cols
                reasons = []
                for c in cols:
                    reasons.extend(parse_reasons(row.get(c)))
                buckets = [bucket_reason(r) for r in reasons]
                for b in buckets:
                    seg_reasons[status][b] += 1
        branch_segment_rec_reasons[branch][seg] = {
            'Yes': reasons_to_top(seg_reasons['Yes']),
            'Maybe': reasons_to_top(seg_reasons['Maybe']),
            'No': reasons_to_top(seg_reasons['No'])
        }

stats['branch_segment_performance'] = branch_segment_perf
stats['branch_segment_recommendation_counts'] = branch_segment_rec_counts
stats['branch_segment_recommendation_reasons'] = branch_segment_rec_reasons

# Fallback post-processing: compute averages from distributions if missing
def weighted_avg_from_distribution(dist_obj):
    if not isinstance(dist_obj, dict):
        return None
    e=g=a=n2=p=0
    for k, v in dist_obj.items():
        try:
            val = int(v) if v is not None else 0
        except Exception:
            continue
        low = str(k).lower()
        lownorm = re.sub(r"[\s./-]", '', low)
        if 'not applicable' in low or 'பொருந்தாது' in low or low in ('na','n/a','n.a') or lownorm=='notapplicable':
            continue
        if 'excellent' in low or 'மிகநன்று' in low or 'மிக நன்று' in low:
            e += val
        elif 'good' in low or 'நன்று' in low:
            g += val
        elif 'average' in low or 'satisfactory' in low or 'சராசரி' in low or 'திருப்தி' in low:
            a += val
        elif 'need' in low or 'needs' in low or 'improve' in low or 'முன்னேற்றம்' in low:
            n2 += val
        elif 'poor' in low or 'மோசம்' in low:
            p += val
    denom = e+g+a+n2+p
    if denom <= 0:
        return None
    num = 5*e + 4*g + 3*a + 2*n2 + 1*p
    return float(num/denom)

# Subjects
for name, info in list(stats.get('subject_performance', {}).items()):
    avg = info.get('average')
    if avg is None or (isinstance(avg, float) and np.isnan(avg)):
        dist = info.get('rating_distribution') or {}
        wa = weighted_avg_from_distribution(dist)
        stats['subject_performance'][name]['average'] = wa

# Category groups
cat_perf = stats.get('category_performance', {})
for grp in list(cat_perf.keys()):
    for name, info in list(cat_perf.get(grp, {}).items()):
        avg = info.get('average')
        if avg is None or (isinstance(avg, float) and np.isnan(avg)):
            dist = info.get('rating_distribution') or {}
            wa = weighted_avg_from_distribution(dist)
            stats['category_performance'][grp][name]['average'] = wa

# Summary category scores fallback
summary = stats.get('summary', {})
def avg_values(d):
    vals = [v for v in d.values() if v is not None and not (isinstance(v, float) and np.isnan(v))]
    return float(sum(vals)/len(vals)) if vals else None

if summary.get('category_scores', {}).get('Academics') in (None,) or np.isnan(summary.get('category_scores', {}).get('Academics')):
    subj_avgs = [v.get('average') for v in (stats.get('subject_performance') or {}).values()]
    val = avg_values({i:a for i,a in enumerate(subj_avgs) if a is not None})
    stats['summary']['category_scores']['Academics'] = val

groups_map = {
    'Environment': 'Environment Quality',
    'Infrastructure': 'Infrastructure',
    'Administration': 'Administrative Support'
}
for key, grp in groups_map.items():
    cs = stats['summary']['category_scores'].get(key)
    if cs is None or (isinstance(cs, float) and np.isnan(cs)):
        items = stats.get('category_performance', {}).get(grp, {})
        val = avg_values({k:items[k].get('average') for k in items})
        stats['summary']['category_scores'][key] = val

# Overall avg fallback
if stats['summary'].get('overall_avg') is None or (isinstance(stats['summary'].get('overall_avg'), float) and np.isnan(stats['summary'].get('overall_avg'))):
    vals = [v for v in (stats['summary'].get('category_scores') or {}).values() if v is not None and not (isinstance(v, float) and np.isnan(v))]
    stats['summary']['overall_avg'] = float(sum(vals)/len(vals)) if vals else None

# Now that all aggregates are computed, clean and save JSON
stats_clean = clean_nan(stats)
with open('/Users/venkubabugollapudi/Desktop/Feedback/Feed Back/feedback_stats.json', 'w', encoding='utf-8') as f:
    json.dump(stats_clean, f, indent=2, ensure_ascii=False)

print("\n✅ Analysis complete! Statistics saved to feedback_stats.json")
print(f"\nTop 3 Branches by Overall Performance:")
for i, (name, score, count) in enumerate(stats['rankings']['branches'][:3], 1):
    print(f"  {i}. {name}: {score:.2f}/5.0 (n={count})")

print(f"\nTop 3 Orientations by Overall Performance:")
for i, (name, score, count) in enumerate(stats['rankings']['orientations'][:3], 1):
    print(f"  {i}. {name}: {score:.2f}/5.0 (n={count})")

# Create a concise PowerPoint presentation
def create_ppt_report(stats_dict, output_file):
    def safe_float(x, default=0.0):
        try:
            return float(x) if x is not None else default
        except Exception:
            return default

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 31, 63)
    # Derive base_dir from first configured input file (XLSX/CSV) if available; fallback to script folder
    try:
        first_input = next((p for p in input_files_used if os.path.exists(p)), None)
    except Exception:
        first_input = None
    base_dir = os.path.dirname(first_input) if first_input else os.path.dirname(__file__)
    # Prefer srichaitanya.jpg, fall back to common names
    for candidate in ['srichaitanya.jpg', 'logo.png', 'srichaitanya.png']:
        logo_path = os.path.join(base_dir, candidate)
        if os.path.exists(logo_path):
            try:
                slide.shapes.add_picture(logo_path, Inches(0.6), Inches(0.6), height=Inches(1.0))
            except Exception:
                pass
            break

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Pre-Primary Sri Chaitanya Techno Schools – Tamil Nadu"
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 215, 0)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Academic & Administrative Review • Academic Year: 2025–26"
    sp = subtitle_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.font.size = Pt(24)
    sp.font.color.rgb = RGBColor(255, 255, 255)

    # Dashboard-style menu (clickable buttons)
    sections = [
        ('exec', 'Executive Summary'),
        ('academic', 'Academic Quality'),
        ('env', 'Environment & Safety'),
        ('comm', 'Communication & Administration'),
        ('infra', 'Infrastructure & Facilities'),
        ('strengths', 'Strengths & Improvements'),
        ('branch', 'Branch Comparison'),
    ]
    section_slides = {}
    menu_slide = prs.slides.add_slide(prs.slide_layouts[6])
    mf = menu_slide.background.fill
    mf.solid()
    mf.fore_color.rgb = RGBColor(240, 248, 255)
    mt = menu_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    mtf = mt.text_frame
    mtf.text = "Dashboard Menu"
    mtp = mtf.paragraphs[0]
    mtp.font.size = Pt(32)
    mtp.font.bold = True
    mtp.font.color.rgb = RGBColor(0, 31, 63)
    menu_buttons = {}
    for idx, (k, label) in enumerate(sections):
        x = 0.8 + (idx % 2) * 4.8
        y = 1.3 + (idx // 2) * 0.9
        btn = menu_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(0.6))
        btn.fill.solid()
        btn.fill.fore_color.rgb = RGBColor(255, 215, 0)
        btn.line.color.rgb = RGBColor(0, 31, 63)
        btf = btn.text_frame
        btf.text = label
        bp = btf.paragraphs[0]
        bp.font.bold = True
        bp.font.size = Pt(18)
        bp.font.color.rgb = RGBColor(0, 31, 63)
        menu_buttons[k] = btn

    # Helper: add persistent tab bar at top of a slide
    def add_tab_bar(slide_obj, active_key=None):
        try:
            start_x = 0.4
            y = 0.05
            tab_w = 1.25
            tab_h = 0.4
            gap = 0.05
            for idx, (key, label) in enumerate(sections):
                x = start_x + idx * (tab_w + gap)
                tab = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(tab_w), Inches(tab_h)
                )
                tab.fill.solid()
                if active_key == key:
                    tab.fill.fore_color.rgb = RGBColor(255, 215, 0)
                else:
                    tab.fill.fore_color.rgb = RGBColor(240, 248, 255)
                tab.line.color.rgb = RGBColor(0, 31, 63)
                tf = tab.text_frame
                tf.text = label.split(' ')[0]
                p = tf.paragraphs[0]
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 31, 63)
                # Link to target slide if available
                target = section_slides.get(key)
                if target is not None:
                    try:
                        tab.click_action.target_slide = target
                    except Exception:
                        pass
            # Menu button at far right
            try:
                menu_btn = slide_obj.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(y), Inches(0.8), Inches(tab_h)
                )
                menu_btn.fill.solid(); menu_btn.fill.fore_color.rgb = RGBColor(255, 215, 0)
                menu_btn.line.color.rgb = RGBColor(0, 31, 63)
                mtf = menu_btn.text_frame; mtf.text = 'Menu'; mtf.paragraphs[0].font.size = Pt(12); mtf.paragraphs[0].font.bold = True
                menu_btn.click_action.target_slide = menu_slide
            except Exception:
                pass
        except Exception:
            pass

    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = tbox.text_frame
    tf.text = "Executive Summary"
    tp = tf.paragraphs[0]
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0, 31, 63)
    # Helper: card container resembling dashboard .chart-container
    def add_card(sl, x, y, w, h, title_text):
        card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(255, 215, 0)
        try:
            card.line.width = Pt(2)
        except Exception:
            pass
        # Title
        tb = sl.shapes.add_textbox(Inches(x+0.3), Inches(y+0.2), Inches(w-0.6), Inches(0.4))
        tft = tb.text_frame; tft.text = title_text
        p0 = tft.paragraphs[0]; p0.font.size = Pt(19); p0.font.bold = True; p0.font.color.rgb = RGBColor(0, 31, 63)
        return card

    # Helper: stat KPI card resembling dashboard .stat-card
    def add_stat_card(sl, x, y, w, h, label, value, rgb):
        box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = rgb
        box.line.color.rgb = RGBColor(255, 255, 255)
        # Label
        lb = sl.shapes.add_textbox(Inches(x+0.2), Inches(y+0.2), Inches(w-0.4), Inches(0.35))
        ltf = lb.text_frame; ltf.text = label
        lp = ltf.paragraphs[0]; lp.font.size = Pt(14); lp.font.bold = True; lp.font.color.rgb = RGBColor(255, 255, 255)
        # Value
        vb = sl.shapes.add_textbox(Inches(x+0.2), Inches(y+0.6), Inches(w-0.4), Inches(0.6))
        vtf = vb.text_frame; vtf.text = str(value)
        vp = vtf.paragraphs[0]; vp.font.size = Pt(28); vp.font.bold = True; vp.font.color.rgb = RGBColor(255, 255, 255)
        return box

    # Helper: style charts to resemble dashboard look
    def style_bar_chart(chart, color=RGBColor(33, 150, 243), show_legend=False, max_scale=5.0):
        try:
            chart.has_legend = show_legend
        except Exception:
            pass
        try:
            # Axes styling
            ca = chart.category_axis
            va = chart.value_axis
            if hasattr(ca, 'has_major_gridlines'):
                ca.has_major_gridlines = False
            if hasattr(va, 'has_major_gridlines'):
                va.has_major_gridlines = False
            if max_scale:
                va.maximum_scale = max_scale
                va.minimum_scale = 0
        except Exception:
            pass
        try:
            for s in chart.series:
                fmt = s.format.fill
                fmt.solid(); fmt.fore_color.rgb = color
        except Exception:
            pass

    def style_pie_chart(chart, show_legend=True):
        try:
            chart.has_legend = show_legend
        except Exception:
            pass

    summary = stats_dict.get('summary', {})
    total_responses = summary.get('total_responses', 0)
    branches_count = len((summary.get('branches') or {}).keys() if isinstance(summary.get('branches'), dict) else summary.get('branches') or [])
    classes_count = len((summary.get('classes') or {}).keys() if isinstance(summary.get('classes'), dict) else summary.get('classes') or [])
    orientations_count = len((summary.get('orientations') or {}).keys() if isinstance(summary.get('orientations'), dict) else summary.get('orientations') or [])

    # KPI stat cards (4 across)
    kx = 0.6; ky = 1.0; kw = 2.2; kh = 1.1; gap = 0.3
    add_stat_card(slide, kx + 0*(kw+gap), ky, kw, kh, 'Total Responses', total_responses, RGBColor(238, 90, 111))
    add_stat_card(slide, kx + 1*(kw+gap), ky, kw, kh, 'Branches', branches_count, RGBColor(68, 160, 141))
    add_stat_card(slide, kx + 2*(kw+gap), ky, kw, kh, 'Classes', classes_count, RGBColor(243, 156, 18))
    add_stat_card(slide, kx + 3*(kw+gap), ky, kw, kh, 'Orientations', orientations_count, RGBColor(52, 31, 151))

    # Left card: Category Scores
    cat_card = add_card(slide, 0.5, 2.4, 4.6, 3.2, 'Overall Satisfaction by Category')
    try:
        cat_scores = (summary.get('category_scores') or {})
        if cat_scores:
            cats = list(cat_scores.keys())
            cs = CategoryChartData(); cs.categories = cats
            cs.add_series('Score', [float(cat_scores.get(k)) if cat_scores.get(k) is not None else 0.0 for k in cats])
            _chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.7), Inches(3.0), Inches(4.2), Inches(2.4), cs
            ).chart
            style_bar_chart(_chart, show_legend=False, max_scale=5.0)
    except Exception:
        pass

    # Right card: Recommendation pie + key reasons
    rec_card = add_card(slide, 5.3, 2.4, 4.2, 3.2, 'Recommendation Status')
    try:
        rec = stats_dict.get('recommendation', {})
        dist = rec.get('distribution') or {}
        if dist:
            cats_rec = list(dist.keys())
            cd = ChartData(); cd.categories = cats_rec
            cd.add_series('Responses', [int(dist.get(k) or 0) for k in cats_rec])
            _pie = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, Inches(5.5), Inches(3.0), Inches(2.2), Inches(2.2), cd
            ).chart
            style_pie_chart(_pie, show_legend=True)
        reasons = stats_dict.get('recommendation_reasons', {})
        yes_top = (reasons.get('Yes') or {}).get('top') or []
        no_top = (reasons.get('No') or {}).get('top') or []
        ylist = 3.0
        ybox = slide.shapes.add_textbox(Inches(7.9), Inches(ylist), Inches(1.4), Inches(0.3))
        ytf = ybox.text_frame; ytf.text = 'Why Yes'; ytf.paragraphs[0].font.bold = True; ytf.paragraphs[0].font.color.rgb = RGBColor(67, 160, 71)
        ylist += 0.35
        for item in yes_top[:3]:
            tb = slide.shapes.add_textbox(Inches(7.9), Inches(ylist), Inches(1.6), Inches(0.3))
            tb.text_frame.text = f"• {item[0][:20]} {item[1]}%"; ylist += 0.3
        nbox = slide.shapes.add_textbox(Inches(7.9), Inches(4.4), Inches(1.4), Inches(0.3))
        ntf = nbox.text_frame; ntf.text = 'Why No'; ntf.paragraphs[0].font.bold = True; ntf.paragraphs[0].font.color.rgb = RGBColor(229, 57, 53)
        ny = 4.75
        for item in no_top[:3]:
            tb = slide.shapes.add_textbox(Inches(7.9), Inches(ny), Inches(1.6), Inches(0.3))
            tb.text_frame.text = f"• {item[0][:20]} {item[1]}%"; ny += 0.3
    except Exception:
        pass

    # Full-width card: Branch-wise Overall
    br_card = add_card(slide, 0.5, 5.8, 9.0, 1.6, 'Branch-wise Overall Rating')
    try:
        top_branches_all = (stats_dict.get('rankings', {}).get('branches') or [])[:10]
        if top_branches_all:
            cdb = CategoryChartData(); cdb.categories = [b[0][:25] for b in top_branches_all]
            cdb.add_series('Overall Score', [safe_float(b[1]) for b in top_branches_all])
            _chart2 = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.7), Inches(6.2), Inches(8.6), Inches(1.0), cdb
            ).chart
            style_bar_chart(_chart2, show_legend=False, max_scale=5.0)
    except Exception:
        pass

    # Back to Menu button
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass
    section_slides['exec'] = slide

    # Top branches chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Top 10 Branches (Overall)"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)

    top_branches = (stats_dict.get('rankings', {}).get('branches') or [])[:10]
    if top_branches:
        chart_data = CategoryChartData()
        chart_data.categories = [b[0][:25] for b in top_branches]
        chart_data.add_series('Overall Score', [safe_float(b[1]) for b in top_branches])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(1), Inches(1.2), Inches(8), Inches(5), chart_data
        ).chart
        style_bar_chart(chart, show_legend=False, max_scale=5.0)
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass
    section_slides['branch'] = slide

    # Branch Recommendation % slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Branch Recommendation %"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)
    br = stats_dict.get('branch_recommendation_pct') or {}
    if br:
        items = sorted(br.items(), key=lambda x: (x[1] is None, -x[1] if x[1] is not None else 0))[:15]
        cd = CategoryChartData()
        cd.categories = [k[:25] for k, _ in items]
        cd.add_series('Yes %', [float(v) if v is not None else 0.0 for _, v in items])
        _br = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(1.0), Inches(1.2), Inches(8), Inches(5), cd
        ).chart
        style_bar_chart(_br, show_legend=False, max_scale=100.0)
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass

    # Academics vs Infrastructure (Scatter)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Academics vs Infrastructure (Branch Scatter)"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)
    try:
        bp = stats_dict.get('branch_performance') or {}
        if bp:
            xy = XyChartData()
            series = xy.add_series('Branches')
            for name, info in bp.items():
                x = safe_float(info.get('subject_avg'))
                y = safe_float(info.get('infrastructure_avg'))
                series.add_data_point(x, y)
            slide.shapes.add_chart(
                XL_CHART_TYPE.XY_SCATTER, Inches(1.0), Inches(1.2), Inches(8), Inches(5), xy
            )
    except Exception:
        pass
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass

    # Subject performance list
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Subject-wise Performance"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)

    subjects = stats_dict.get('subject_performance', {})
    # Card container and bar chart for subject-wise performance
    _ = add_card(slide, 0.5, 1.0, 9.0, 5.8, 'Subject-wise Performance')
    if subjects:
        cats = list(subjects.keys())
        data = [safe_float(subjects[k].get('average')) for k in cats]
        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series('Average', data)
        ch = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.0), cd
        ).chart
        style_bar_chart(ch, show_legend=False, max_scale=5.0)
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass
    section_slides['academic'] = slide

    # Teaching Indicators & PTM slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Teaching Indicators & PTM"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)
    # Teaching indicators bar
    ti = stats_dict.get('teaching_indicators') or {}
    if ti:
        cats = list(ti.keys())
        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series('Avg', [safe_float(ti[k]) for k in cats])
        _ti = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.2), Inches(5.0), Inches(4.8), cd
        ).chart
        style_bar_chart(_ti, show_legend=False, max_scale=5.0)
    # PTM effectiveness donut
    ptm = stats_dict.get('ptm_effectiveness')
    try:
        if ptm is not None:
            val = max(0.0, min(5.0, safe_float(ptm)))
            cd2 = ChartData()
            cd2.categories = ['Effective', 'Gap']
            cd2.add_series('PTM', [val, max(0.0, 5.0 - val)])
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.DOUGHNUT, Inches(6.3), Inches(1.5), Inches(3.8), Inches(3.8), cd2
            ).chart
            chart.has_legend = True
    except Exception:
        pass
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass

    # Environment metrics (top 4)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Environment Quality (Top Metrics)"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)

    env = stats_dict.get('category_performance', {}).get('Environment Quality', {})
    top_env = sorted(env.items(), key=lambda x: safe_float(x[1].get('average')), reverse=True)[:4]
    y = 1.3
    for metric, info in top_env:
        sb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(6.5), Inches(0.5))
        stf = sb.text_frame
        stf.text = metric[:60]
        sp = stf.paragraphs[0]
        sp.font.size = Pt(18)
        sp.font.color.rgb = RGBColor(0, 31, 63)

        vb = slide.shapes.add_textbox(Inches(7.4), Inches(y), Inches(1.6), Inches(0.5))
        vtf = vb.text_frame
        vtf.text = f"{safe_float(info.get('average')):.2f}/5.0"
        vp = vtf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vp.font.size = Pt(18)
        vp.font.bold = True
        vp.font.color.rgb = RGBColor(76, 175, 80)
        y += 0.6
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass
    section_slides['env'] = slide

    # Communication & Administration slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Communication & Administration"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)
    # Communication metrics bar and Admin support bar (styled)
    comm = stats_dict.get('communication_metrics') or {}
    if comm:
        cats = list(comm.keys())[:8]
        cd = CategoryChartData(); cd.categories = cats
        cd.add_series('Avg', [safe_float(comm[k]) for k in cats])
        _cm = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.2), Inches(4.2), Inches(4.5), cd
        ).chart
        style_bar_chart(_cm, show_legend=False, max_scale=5.0)
    # Admin support bar
    admin = (stats_dict.get('category_performance', {}).get('Administrative Support') or {})
    if admin:
        items = list(admin.items())[:8]
        cd2 = CategoryChartData()
        cd2.categories = [k[:25] for k, _ in items]
        cd2.add_series('Avg', [safe_float(v.get('average')) for _, v in items])
        _adm = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(5.2), Inches(1.2), Inches(4.2), Inches(4.5), cd2
        ).chart
        style_bar_chart(_adm, show_legend=False, max_scale=5.0)
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass
    section_slides['comm'] = slide

    # Concern Handling & Resolution slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Concern Handling & Resolution"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)
    # Concern roles bar
    roles = stats_dict.get('concern_roles') or {}
    if roles:
        cats = list(roles.keys())
        cd = CategoryChartData(); cd.categories = cats
        cd.add_series('Avg', [safe_float(roles[k]) for k in cats])
        _roles = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.2), Inches(4.2), Inches(4.5), cd
        ).chart
        style_bar_chart(_roles, show_legend=False, max_scale=5.0)
    # Concern resolution pie
    cres = stats_dict.get('concern_resolution') or {}
    if cres:
        cats_cres = list(cres.keys())
        cd2 = ChartData(); cd2.categories = cats_cres
        cd2.add_series('Responses', [int(cres.get(k) or 0) for k in cats_cres])
        _cres = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, Inches(5.2), Inches(1.2), Inches(4.2), Inches(4.5), cd2
        ).chart
        style_pie_chart(_cres, show_legend=True)
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass

    # Infrastructure slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Infrastructure & Facilities"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)
    infra = (stats_dict.get('category_performance', {}).get('Infrastructure') or {})
    if infra:
        items = sorted(infra.items(), key=lambda x: safe_float(x[1].get('average')), reverse=True)[:10]
        cd = CategoryChartData()
        cd.categories = [k[:25] for k, _ in items]
        cd.add_series('Avg', [safe_float(v.get('average')) for _, v in items])
        _infra = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(1.0), Inches(1.2), Inches(8), Inches(5), cd
        ).chart
        style_bar_chart(_infra, show_legend=False, max_scale=5.0)
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass
    section_slides['infra'] = slide

    # Strengths & Improvements slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tr = t.text_frame
    tr.text = "Strengths & Improvements"
    tr.paragraphs[0].font.size = Pt(28)
    tr.paragraphs[0].font.bold = True
    tr.paragraphs[0].font.color.rgb = RGBColor(0, 31, 63)
    reasons = stats_dict.get('recommendation_reasons', {})
    yes_top = (reasons.get('Yes') or {}).get('top') or []
    no_top = (reasons.get('No') or {}).get('top') or []
    ybox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4.2), Inches(0.5))
    ytf = ybox.text_frame; ytf.text = 'Top Strengths (Why Yes)'; ytf.paragraphs[0].font.bold = True
    yy = 1.7
    for item in yes_top[:8]:
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(yy), Inches(3.8), Inches(0.4))
        tfb = tb.text_frame; tfb.text = f"• {item[0]} - {item[1]}%"; yy += 0.35
    nbox = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.2), Inches(0.5))
    ntf = nbox.text_frame; ntf.text = 'Top Improvements (Why No)'; ntf.paragraphs[0].font.bold = True
    ny = 1.7
    for item in no_top[:8]:
        tb = slide.shapes.add_textbox(Inches(5.4), Inches(ny), Inches(3.8), Inches(0.4))
        tfb = tb.text_frame; tfb.text = f"• {item[0]} - {item[1]}%"; ny += 0.35
    try:
        back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.2), Inches(1.2), Inches(0.5))
        back.fill.solid(); back.fill.fore_color.rgb = RGBColor(255, 215, 0)
        back.line.color.rgb = RGBColor(0, 31, 63)
        bt = back.text_frame; bt.text = 'Menu'; bt.paragraphs[0].font.bold = True
        back.click_action.target_slide = menu_slide
    except Exception:
        pass
    section_slides['strengths'] = slide

    # Wire menu buttons to their respective slides
    try:
        for key, btn in (menu_buttons or {}).items():
            target = section_slides.get(key)
            if target is not None:
                try:
                    btn.click_action.target_slide = target
                except Exception:
                    pass
    except Exception:
        pass

    # Add tab bars to all slides, highlighting the active section slide when applicable
    try:
        by_id = {id(v): k for k, v in section_slides.items()}
        for s in prs.slides:
            active = by_id.get(id(s))
            add_tab_bar(s, active)
    except Exception:
        pass

    prs.save(output_file)

# Generate PPT next to JSON
ppt_path = '/Users/venkubabugollapudi/Desktop/Feedback/Feed Back/Pre_Primary_Feedback_Analysis.pptx'
create_ppt_report(stats_clean, ppt_path)
print(f"\n📁 PowerPoint saved: {ppt_path}")
